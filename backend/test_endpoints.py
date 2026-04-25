import io
import json
import unittest
import zipfile
from unittest.mock import patch

from docx import Document
from fastapi.testclient import TestClient
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from PyPDF2 import PdfReader
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from main import app


def build_pdf_bytes(page_count: int = 3) -> bytes:
    buffer = io.BytesIO()
    pdf_canvas = canvas.Canvas(buffer, pagesize=letter)

    for page_number in range(1, page_count + 1):
        pdf_canvas.setFont("Helvetica", 16)
        pdf_canvas.drawString(72, 720, f"Test Page {page_number}")
        pdf_canvas.setFont("Helvetica", 12)
        pdf_canvas.drawString(72, 690, "NOVA PDF integration test")
        pdf_canvas.showPage()

    pdf_canvas.save()
    buffer.seek(0)
    return buffer.read()


def build_sensitive_pdf_bytes() -> bytes:
    buffer = io.BytesIO()
    pdf_canvas = canvas.Canvas(buffer, pagesize=letter)
    pdf_canvas.setFont("Helvetica", 14)
    pdf_canvas.drawString(72, 720, "Dossier SECRET du client NOVA")
    pdf_canvas.drawString(72, 700, "Ne pas diffuser.")
    pdf_canvas.save()
    buffer.seek(0)
    return buffer.read()


def build_png_bytes() -> bytes:
    image = Image.new("RGB", (320, 200), color=(240, 170, 120))
    output = io.BytesIO()
    image.save(output, format="PNG")
    output.seek(0)
    return output.read()


def build_docx_bytes() -> bytes:
    document = Document()
    document.add_heading("Demo DOCX", level=1)
    document.add_paragraph("Bonjour depuis Word.")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "A1"
    table.rows[0].cells[1].text = "B1"
    table.rows[1].cells[0].text = "A2"
    table.rows[1].cells[1].text = "B2"
    output = io.BytesIO()
    document.save(output)
    output.seek(0)
    return output.read()


def build_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    sheet["A1"] = "Produit"
    sheet["B1"] = "Prix"
    sheet["A2"] = "PDF"
    sheet["B2"] = 19
    output = io.BytesIO()
    workbook.save(output)
    output.seek(0)
    return output.read()


def build_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide demo"
    slide.placeholders[1].text = "Bonjour depuis PowerPoint."
    output = io.BytesIO()
    presentation.save(output)
    output.seek(0)
    return output.read()


def pdf_reader_from_bytes(content: bytes) -> PdfReader:
    return PdfReader(io.BytesIO(content))


class NovaPdfEndpointsTest(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.client = TestClient(app)
        cls.sample_pdf = build_pdf_bytes(3)
        cls.second_pdf = build_pdf_bytes(2)
        cls.sample_png = build_png_bytes()
        cls.sensitive_pdf = build_sensitive_pdf_bytes()
        cls.sample_docx = build_docx_bytes()
        cls.sample_xlsx = build_xlsx_bytes()
        cls.sample_pptx = build_pptx_bytes()
        cls.sample_html = b"<html><body><h1>NOVA</h1><p>Bonjour HTML.</p></body></html>"

    def test_homepage_renders(self):
        response = self.client.get("/")
        self.assertEqual(response.status_code, 200)
        self.assertIn("NOVA PDF Suite", response.text)
        self.assertIn("Convertir depuis PDF", response.text)

    def test_merge_endpoint(self):
        response = self.client.post(
            "/api/merge",
            files=[
                ("files", ("first.pdf", self.sample_pdf, "application/pdf")),
                ("files", ("second.pdf", self.second_pdf, "application/pdf")),
            ],
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(len(pdf_reader_from_bytes(response.content).pages), 5)

    def test_reorder_endpoint(self):
        response = self.client.post(
            "/api/reorder",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
            data={"pages": "3,1"},
        )
        self.assertEqual(response.status_code, 200)
        self.assertIn("Test Page 3", pdf_reader_from_bytes(response.content).pages[0].extract_text())

    def test_reorder_range_includes_endpoints(self):
        response = self.client.post(
            "/api/reorder",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
            data={"pages": "3-1"},
        )
        self.assertEqual(response.status_code, 200)
        reader = pdf_reader_from_bytes(response.content)
        texts = [(page.extract_text() or "") for page in reader.pages]
        self.assertEqual(len(texts), 3)
        self.assertIn("Test Page 3", texts[0])
        self.assertIn("Test Page 2", texts[1])
        self.assertIn("Test Page 1", texts[2])

    def test_crop_endpoint(self):
        response = self.client.post(
            "/api/crop",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
            data={"top": "20", "right": "20", "bottom": "20", "left": "20"},
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(len(pdf_reader_from_bytes(response.content).pages), 3)

    def test_watermark_endpoint(self):
        response = self.client.post(
            "/api/watermark",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
            data={"text": "CONFIDENTIEL", "opacity": "0.5"},
        )
        self.assertEqual(response.status_code, 200)
        self.assertIn("CONFIDENTIEL", pdf_reader_from_bytes(response.content).pages[0].extract_text())

    def test_pdf_to_jpg_endpoint(self):
        response = self.client.post(
            "/api/pdf-to-jpg",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
        )
        self.assertEqual(response.status_code, 200)
        archive = zipfile.ZipFile(io.BytesIO(response.content))
        self.assertEqual(sorted(archive.namelist()), ["page-1.jpg", "page-2.jpg", "page-3.jpg"])

    def test_pdf_to_word_endpoint(self):
        response = self.client.post(
            "/api/pdf-to-word",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
        )
        self.assertEqual(response.status_code, 200)
        self.assertTrue(response.content.startswith(b"PK"))

    def test_pdf_to_excel_endpoint(self):
        response = self.client.post(
            "/api/pdf-to-excel",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
        )
        self.assertEqual(response.status_code, 200)
        self.assertTrue(response.content.startswith(b"PK"))

    def test_numbering_endpoint(self):
        response = self.client.post(
            "/api/numbering",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
            data={"format_str": "Page {page}/{total}", "position": "bottom-center"},
        )
        self.assertEqual(response.status_code, 200)
        self.assertIn("Page 1/3", pdf_reader_from_bytes(response.content).pages[0].extract_text())

    def test_repair_endpoint(self):
        response = self.client.post(
            "/api/repair",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(len(pdf_reader_from_bytes(response.content).pages), 3)

    def test_image_to_pdf_endpoint(self):
        response = self.client.post(
            "/api/convert/image-to-pdf",
            files={"file": ("sample.png", self.sample_png, "image/png")},
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(len(pdf_reader_from_bytes(response.content).pages), 1)

    def test_html_to_pdf_endpoint(self):
        response = self.client.post(
            "/api/convert/html-to-pdf",
            files={"file": ("sample.html", self.sample_html, "text/html")},
        )
        self.assertEqual(response.status_code, 200)
        self.assertGreater(len(response.content), 500)

    def test_word_to_pdf_endpoint(self):
        response = self.client.post(
            "/api/convert/word-to-pdf",
            files={"file": ("sample.docx", self.sample_docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
        )
        self.assertEqual(response.status_code, 200)
        self.assertGreater(len(response.content), 500)

    def test_excel_to_pdf_endpoint(self):
        response = self.client.post(
            "/api/convert/excel-to-pdf",
            files={"file": ("sample.xlsx", self.sample_xlsx, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
        )
        self.assertEqual(response.status_code, 200)
        self.assertGreater(len(response.content), 500)

    def test_powerpoint_to_pdf_endpoint(self):
        response = self.client.post(
            "/api/convert/powerpoint-to-pdf",
            files={"file": ("sample.pptx", self.sample_pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        self.assertEqual(response.status_code, 200)
        self.assertGreater(len(response.content), 500)

    def test_protect_and_unlock_endpoints(self):
        protected = self.client.post(
            "/api/protect",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
            data={"user_password": "secret123"},
        )
        self.assertEqual(protected.status_code, 200)
        self.assertTrue(pdf_reader_from_bytes(protected.content).is_encrypted)

        unlocked = self.client.post(
            "/api/unlock",
            files={"file": ("protected.pdf", protected.content, "application/pdf")},
            data={"password": "secret123"},
        )
        self.assertEqual(unlocked.status_code, 200)
        self.assertFalse(pdf_reader_from_bytes(unlocked.content).is_encrypted)

    def test_compare_endpoint(self):
        response = self.client.post(
            "/api/compare",
            files=[
                ("file_a", ("first.pdf", self.sample_pdf, "application/pdf")),
                ("file_b", ("second.pdf", self.second_pdf, "application/pdf")),
            ],
        )
        self.assertEqual(response.status_code, 200)
        report = json.loads(response.content.decode("utf-8"))
        self.assertFalse(report["same_page_count"])

    def test_censor_endpoint(self):
        response = self.client.post(
            "/api/censor",
            files={"file": ("sensitive.pdf", self.sensitive_pdf, "application/pdf")},
            data={"terms": "SECRET"},
        )
        self.assertEqual(response.status_code, 200)
        redacted_text = pdf_reader_from_bytes(response.content).pages[0].extract_text() or ""
        self.assertNotIn("SECRET", redacted_text)

    def test_sign_endpoint(self):
        response = self.client.post(
            "/api/sign",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
            data={"signer_name": "NOVA Rescue", "reason": "Validation", "location": "Lab"},
        )
        self.assertEqual(response.status_code, 200)
        first_page_text = pdf_reader_from_bytes(response.content).pages[0].extract_text() or ""
        self.assertIn("Signe par NOVA Rescue", first_page_text)
        self.assertIn("Signature visible non cryptographique", first_page_text)

    def test_summarize_endpoint(self):
        with patch("pdf_utils._call_openai_text", return_value=None):
            response = self.client.post(
                "/api/ai/summarize",
                files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
                data={"max_sentences": "3"},
            )
        self.assertEqual(response.status_code, 200)
        self.assertIn("Test Page", response.content.decode("utf-8"))

    def test_translate_endpoint(self):
        with patch("pdf_utils._call_openai_text", return_value=None):
            response = self.client.post(
                "/api/ai/translate",
                files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
                data={"target_language": "english"},
            )
        self.assertEqual(response.status_code, 200)
        text = "\n".join((page.extract_text() or "") for page in pdf_reader_from_bytes(response.content).pages)
        self.assertIn("Traduction locale vers english", text)

    def test_invalid_page_range_returns_400(self):
        response = self.client.post(
            "/api/split",
            files={"file": ("sample.pdf", self.sample_pdf, "application/pdf")},
            data={"pages": "99"},
        )
        self.assertEqual(response.status_code, 400)
        self.assertIn("n'existe pas", response.json()["detail"])

    def test_upload_size_limit_returns_400(self):
        with patch("pdf_utils.MAX_UPLOAD_BYTES", 32):
            response = self.client.post(
                "/api/compress",
                files={"file": ("oversized.pdf", b"x" * 33, "application/pdf")},
            )
        self.assertEqual(response.status_code, 400)
        self.assertIn("depasse la limite", response.json()["detail"])


if __name__ == "__main__":
    unittest.main(verbosity=2)
