import io
import json
import os
import subprocess
import shutil
import tempfile
import textwrap
import urllib.error
import urllib.request
import zipfile
from datetime import datetime, timezone
from html.parser import HTMLParser
from pathlib import Path
from typing import Iterable, List, Optional

import pdfplumber
import pikepdf
import pypdfium2 as pdfium
from docx import Document
from openpyxl import Workbook, load_workbook
from PIL import Image, ImageDraw, UnidentifiedImageError
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader, PdfWriter
from PyPDF2.errors import PdfReadError
from reportlab.lib import colors
from pdf2image import convert_from_path
from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError, PDFSyntaxError
from reportlab.lib.colors import black
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfgen import canvas

TEMP_SUFFIXES = {
    ".docx",
    ".gif",
    ".htm",
    ".html",
    ".jpeg",
    ".jpg",
    ".pdf",
    ".png",
    ".pptx",
    ".tiff",
    ".txt",
    ".webp",
    ".xlsx",
}
PAGE_NUMBER_POSITIONS = {
    "top-left",
    "top-center",
    "top-right",
    "bottom-left",
    "bottom-center",
    "bottom-right",
}
MAX_UPLOAD_BYTES = int(os.getenv("NOVA_MAX_UPLOAD_BYTES", str(50 * 1024 * 1024)))
COMMAND_TIMEOUT_SECONDS = int(os.getenv("NOVA_COMMAND_TIMEOUT_SECONDS", "120"))
COPY_CHUNK_BYTES = 1024 * 1024
AI_MAX_INPUT_CHARS = int(os.getenv("NOVA_AI_MAX_INPUT_CHARS", "24000"))
AI_TIMEOUT_SECONDS = int(os.getenv("NOVA_AI_TIMEOUT_SECONDS", "45"))
AI_MODEL = os.getenv("NOVA_AI_MODEL", "gpt-4o-mini")
SIGNATURE_POSITIONS = {
    "bottom-left",
    "bottom-right",
    "top-left",
    "top-right",
}


class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: List[str] = []

    def handle_starttag(self, tag, attrs):
        if tag in {"p", "br", "div", "section", "article", "li", "tr", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")

    def handle_data(self, data):
        text = data.strip()
        if text:
            self.parts.append(text)

    def get_text(self) -> str:
        joined = " ".join(self.parts)
        lines = [line.strip() for line in joined.splitlines()]
        return "\n".join(line for line in lines if line)


def _find_command(*names: str) -> Optional[str]:
    for name in names:
        path = shutil.which(name)
        if path:
            return path
    return None


def _run_command(command: List[str], error_message: str):
    try:
        subprocess.run(
            command,
            check=True,
            capture_output=True,
            text=True,
            timeout=COMMAND_TIMEOUT_SECONDS,
        )
    except FileNotFoundError as exc:
        raise ValueError(error_message) from exc
    except subprocess.TimeoutExpired as exc:
        raise ValueError(f"{error_message} Delai d'execution depasse.") from exc
    except subprocess.CalledProcessError as exc:
        detail = (exc.stderr or exc.stdout or "").strip()
        if detail:
            raise ValueError(f"{error_message} Detail: {detail}") from exc
        raise ValueError(error_message) from exc


def _convert_with_libreoffice(source_path: Path) -> Optional[bytes]:
    soffice = _find_command("soffice", "libreoffice")
    if not soffice:
        return None

    output_dir = Path(tempfile.mkdtemp())
    try:
        try:
            command = [
                soffice,
                "--headless",
                "--convert-to",
                "pdf:writer_pdf_Export",
                "--outdir",
                str(output_dir),
                str(source_path),
            ]
            _run_command(command, "LibreOffice n'a pas pu convertir ce document en PDF.")
            converted_path = output_dir / f"{source_path.stem}.pdf"
            if converted_path.exists():
                return converted_path.read_bytes()
        except ValueError:
            return None
        return None
    finally:
        shutil.rmtree(output_dir, ignore_errors=True)


def _ocr_with_ocrmypdf(source_path: Path, languages: str) -> Optional[bytes]:
    ocrmypdf = _find_command("ocrmypdf")
    if not ocrmypdf:
        return None

    output_pdf = Path(tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name)
    sidecar_txt = Path(tempfile.NamedTemporaryFile(delete=False, suffix=".txt").name)
    try:
        try:
            command = [
                ocrmypdf,
                "--force-ocr",
                "--skip-big",
                "50",
                "--sidecar",
                str(sidecar_txt),
                "--language",
                languages,
                str(source_path),
                str(output_pdf),
            ]
            _run_command(command, "OCRmyPDF n'a pas pu traiter ce document.")
            if sidecar_txt.exists():
                text = sidecar_txt.read_text(encoding="utf-8", errors="ignore").strip()
                if text:
                    return text.encode("utf-8")
            if output_pdf.exists():
                page_texts = _extract_pdf_text_by_page(output_pdf)
                if any(page_texts):
                    return "\n".join(page_texts).encode("utf-8")
        except ValueError:
            return None
        return b""
    finally:
        cleanup([output_pdf, sidecar_txt])


def save_upload_file(upload_file) -> Path:
    if not getattr(upload_file, "filename", None):
        raise ValueError("Aucun fichier n'a ete fourni.")

    suffix = Path(upload_file.filename).suffix.lower() or ".pdf"
    if suffix not in TEMP_SUFFIXES:
        suffix = ".pdf"

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    upload_file.file.seek(0)
    with tmp:
        copied = 0
        while True:
            chunk = upload_file.file.read(COPY_CHUNK_BYTES)
            if not chunk:
                break
            copied += len(chunk)
            if copied > MAX_UPLOAD_BYTES:
                tmp.close()
                Path(tmp.name).unlink(missing_ok=True)
                raise ValueError(
                    f"Le fichier depasse la limite autorisee de {MAX_UPLOAD_BYTES // (1024 * 1024)} Mo."
                )
            tmp.write(chunk)
    return Path(tmp.name)


def cleanup(paths: Iterable[Path]):
    for path in paths:
        try:
            path.unlink()
        except FileNotFoundError:
            pass
        except Exception:
            pass


def _load_pdf_reader(source_path: Path) -> PdfReader:
    try:
        return PdfReader(str(source_path))
    except PdfReadError as exc:
        raise ValueError("Le fichier PDF est invalide ou corrompu.") from exc


def _serialize_writer(writer: PdfWriter) -> bytes:
    buffer = io.BytesIO()
    writer.write(buffer)
    buffer.seek(0)
    return buffer.read()


def _dedupe_preserve_order(pages: Iterable[int]) -> List[int]:
    ordered = []
    seen = set()
    for page_index in pages:
        if page_index not in seen:
            seen.add(page_index)
            ordered.append(page_index)
    return ordered


def parse_page_ranges(ranges: str, total_pages: int) -> List[int]:
    if total_pages < 1:
        raise ValueError("Le document PDF ne contient aucune page.")
    if not ranges or not ranges.strip():
        raise ValueError("Indiquez au moins une page ou une plage de pages.")

    selected_pages = []
    for part in ranges.split(","):
        token = part.strip()
        if not token:
            continue
        if "-" in token:
            if token.count("-") != 1:
                raise ValueError(f"Plage invalide: '{token}'.")
            start_text, end_text = token.split("-", 1)
            if not start_text.strip().isdigit() or not end_text.strip().isdigit():
                raise ValueError(f"Plage invalide: '{token}'.")
            start_page = int(start_text.strip())
            end_page = int(end_text.strip())
            if start_page < 1 or end_page < 1:
                raise ValueError("Les numeros de page doivent etre superieurs a 0.")
            if start_page > end_page:
                raise ValueError(f"Plage invalide: '{token}'.")
            if end_page > total_pages:
                raise ValueError(f"La page {end_page} n'existe pas. Le document contient {total_pages} pages.")
            selected_pages.extend(range(start_page - 1, end_page))
            continue

        if not token.isdigit():
            raise ValueError(f"Numero de page invalide: '{token}'.")
        page_number = int(token)
        if page_number < 1:
            raise ValueError("Les numeros de page doivent etre superieurs a 0.")
        if page_number > total_pages:
            raise ValueError(f"La page {page_number} n'existe pas. Le document contient {total_pages} pages.")
        selected_pages.append(page_number - 1)

    selected_pages = _dedupe_preserve_order(selected_pages)
    if not selected_pages:
        raise ValueError("Aucune page valide n'a ete selectionnee.")
    return selected_pages


def parse_page_sequence(ranges: str, total_pages: int) -> List[int]:
    if not ranges or not ranges.strip():
        raise ValueError("Indiquez l'ordre des pages.")

    ordered_pages = []
    for part in ranges.split(","):
        token = part.strip()
        if not token:
            continue
        if "-" in token:
            if token.count("-") != 1:
                raise ValueError(f"Sequence invalide: '{token}'.")
            start_text, end_text = token.split("-", 1)
            if not start_text.strip().isdigit() or not end_text.strip().isdigit():
                raise ValueError(f"Sequence invalide: '{token}'.")
            start_page = int(start_text.strip())
            end_page = int(end_text.strip())
            if start_page < 1 or end_page < 1 or start_page > total_pages or end_page > total_pages:
                raise ValueError("L'ordre des pages contient une page inexistante.")
            step = 1 if start_page <= end_page else -1
            ordered_pages.extend(page_number - 1 for page_number in range(start_page, end_page + step, step))
            continue

        if not token.isdigit():
            raise ValueError(f"Numero de page invalide: '{token}'.")
        page_number = int(token)
        if page_number < 1 or page_number > total_pages:
            raise ValueError("L'ordre des pages contient une page inexistante.")
        ordered_pages.append(page_number - 1)

    if not ordered_pages:
        raise ValueError("L'ordre des pages est vide.")
    return ordered_pages


def _page_dimensions(page) -> tuple[float, float]:
    return float(page.mediabox.width), float(page.mediabox.height)


def _build_overlay_page(width: float, height: float, drawer):
    buffer = io.BytesIO()
    pdf_canvas = canvas.Canvas(buffer, pagesize=(width, height))
    drawer(pdf_canvas)
    pdf_canvas.save()
    buffer.seek(0)
    return PdfReader(buffer).pages[0]


def _fit_font_size(text: str, font_name: str, max_width: float, start_size: float, min_size: float) -> float:
    font_size = start_size
    while font_size > min_size and pdfmetrics.stringWidth(text, font_name, font_size) > max_width:
        font_size -= 2
    return max(font_size, min_size)


def _page_number_coordinates(width: float, height: float, position: str) -> tuple[float, float, str]:
    vertical, horizontal = position.split("-")
    margin_x = max(width * 0.05, 24)
    margin_y = max(height * 0.04, 24)

    if horizontal == "left":
        x = margin_x
    elif horizontal == "center":
        x = width / 2
    else:
        x = width - margin_x

    y = height - margin_y if vertical == "top" else margin_y
    return max(x, 8), max(y, 8), horizontal


def _draw_watermark(pdf_canvas, width: float, height: float, text: str, font_size: float, opacity: float):
    pdf_canvas.saveState()
    try:
        pdf_canvas.setFillAlpha(opacity)
    except AttributeError:
        pass
    pdf_canvas.setFillColorRGB(0.42, 0.42, 0.42)
    pdf_canvas.setFont("Helvetica-Bold", font_size)
    pdf_canvas.translate(width / 2, height / 2)
    pdf_canvas.rotate(45)
    pdf_canvas.drawCentredString(0, 0, text)
    pdf_canvas.restoreState()


def _draw_page_number(pdf_canvas, text: str, x: float, y: float, font_name: str, font_size: float, align: str):
    pdf_canvas.saveState()
    pdf_canvas.setFont(font_name, font_size)
    pdf_canvas.setFillColorRGB(0.15, 0.15, 0.15)
    if align == "right":
        pdf_canvas.drawRightString(x, y, text)
    elif align == "center":
        pdf_canvas.drawCentredString(x, y, text)
    else:
        pdf_canvas.drawString(x, y, text)
    pdf_canvas.restoreState()


def _signature_box_coordinates(width: float, height: float, position: str) -> tuple[float, float, float, float]:
    box_width = min(max(width * 0.36, 190), width - 48)
    box_height = 82
    margin = 28
    if position.endswith("right"):
        x = width - box_width - margin
    else:
        x = margin
    if position.startswith("top"):
        y = height - box_height - margin
    else:
        y = margin
    return x, y, box_width, box_height


def _draw_signature_box(pdf_canvas, width: float, height: float, lines: List[str], position: str):
    x, y, box_width, box_height = _signature_box_coordinates(width, height, position)
    pdf_canvas.saveState()
    pdf_canvas.setStrokeColor(colors.HexColor("#1f2937"))
    pdf_canvas.setFillColor(colors.Color(1, 1, 1, alpha=0.88))
    pdf_canvas.roundRect(x, y, box_width, box_height, 8, stroke=1, fill=1)
    pdf_canvas.setFillColor(colors.HexColor("#111827"))
    pdf_canvas.setFont("Helvetica-Bold", 10)
    pdf_canvas.drawString(x + 12, y + box_height - 22, lines[0])
    pdf_canvas.setFont("Helvetica", 8)
    current_y = y + box_height - 38
    for line in lines[1:]:
        pdf_canvas.drawString(x + 12, current_y, line[:72])
        current_y -= 11
    pdf_canvas.restoreState()


def _render_pdf_with_pdfium(source_path: Path) -> List[Image.Image]:
    document = None
    rendered_images: List[Image.Image] = []
    try:
        document = pdfium.PdfDocument(str(source_path))
        scale = 200 / 72
        for page_index in range(len(document)):
            page = document[page_index]
            bitmap = page.render(scale=scale)
            image = bitmap.to_pil().convert("RGB")
            rendered_images.append(image)
            bitmap.close()
            page.close()
        return rendered_images
    except pdfium.PdfiumError as exc:
        raise ValueError("Impossible de convertir ce PDF en images.") from exc
    finally:
        if document is not None:
            document.close()


def _render_pdf_pages(source_path: Path) -> List[Image.Image]:
    try:
        return [image.convert("RGB") for image in convert_from_path(str(source_path), dpi=200, fmt="jpeg")]
    except (PDFInfoNotInstalledError, FileNotFoundError, OSError):
        return _render_pdf_with_pdfium(source_path)
    except (PDFPageCountError, PDFSyntaxError) as exc:
        raise ValueError("Impossible de lire les pages de ce PDF.") from exc


def _extract_pdf_text_by_page(source_path: Path) -> List[str]:
    texts: List[str] = []
    try:
        with pdfplumber.open(str(source_path)) as pdf:
            for page in pdf.pages:
                text = (page.extract_text() or "").strip()
                texts.append(text)
    except Exception:
        reader = _load_pdf_reader(source_path)
        for page in reader.pages:
            texts.append((page.extract_text() or "").strip())
    return texts


def _normalize_lines(lines: Iterable[str]) -> List[str]:
    normalized = []
    for line in lines:
        clean = str(line).replace("\t", " ").strip()
        if clean:
            normalized.append(clean)
    return normalized or ["Document vide."]


def _text_lines_to_pdf(lines: Iterable[str], title: Optional[str] = None) -> bytes:
    normalized = _normalize_lines(lines)
    buffer = io.BytesIO()
    pdf_canvas = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4
    y = height - 48

    if title:
        pdf_canvas.setFont("Helvetica-Bold", 16)
        pdf_canvas.drawString(40, y, title)
        y -= 28

    pdf_canvas.setFont("Helvetica", 10)
    for line in normalized:
        words = line.split()
        current = ""
        for word in words:
            candidate = f"{current} {word}".strip()
            if pdfmetrics.stringWidth(candidate, "Helvetica", 10) > width - 80:
                pdf_canvas.drawString(40, y, current)
                y -= 14
                current = word
                if y < 48:
                    pdf_canvas.showPage()
                    pdf_canvas.setFont("Helvetica", 10)
                    y = height - 48
            else:
                current = candidate
        if current:
            pdf_canvas.drawString(40, y, current)
            y -= 14
        if y < 48:
            pdf_canvas.showPage()
            pdf_canvas.setFont("Helvetica", 10)
            y = height - 48

    pdf_canvas.save()
    buffer.seek(0)
    return buffer.read()


def _extract_full_pdf_text(source_path: Path) -> str:
    page_texts = _extract_pdf_text_by_page(source_path)
    text = "\n\n".join(text for text in page_texts if text.strip()).strip()
    if not text:
        raise ValueError("Aucun texte exploitable trouve dans ce PDF.")
    return text


def _call_openai_text(system_prompt: str, user_prompt: str) -> Optional[str]:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        return None

    payload = {
        "model": AI_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt[:AI_MAX_INPUT_CHARS]},
        ],
        "temperature": 0.2,
    }
    request = urllib.request.Request(
        "https://api.openai.com/v1/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=AI_TIMEOUT_SECONDS) as response:
            body = json.loads(response.read().decode("utf-8"))
    except (urllib.error.URLError, TimeoutError, json.JSONDecodeError, KeyError):
        return None
    choices = body.get("choices") or []
    if not choices:
        return None
    content = choices[0].get("message", {}).get("content", "")
    return content.strip() or None


def _sentences_from_text(text: str) -> List[str]:
    normalized = " ".join(text.replace("\r", " ").replace("\n", " ").split())
    sentences: List[str] = []
    current = []
    for char in normalized:
        current.append(char)
        if char in ".!?":
            sentence = "".join(current).strip()
            if len(sentence) > 24:
                sentences.append(sentence)
            current = []
    tail = "".join(current).strip()
    if len(tail) > 24:
        sentences.append(tail)
    return sentences


def _local_summary(text: str, max_sentences: int) -> str:
    sentences = _sentences_from_text(text)
    if not sentences:
        return textwrap.shorten(text, width=1200, placeholder="...")
    selected = sentences[: max(1, min(max_sentences, len(sentences)))]
    bullets = "\n".join(f"- {sentence}" for sentence in selected)
    return f"Resume local\n\n{bullets}"


def _local_translate_text(text: str, target_language: str) -> str:
    language = target_language.strip() or "francais"
    replacements = {
        "english": {
            "bonjour": "hello",
            "page": "page",
            "document": "document",
            "confidentiel": "confidential",
            "client": "client",
            "prix": "price",
            "produit": "product",
        },
        "francais": {
            "hello": "bonjour",
            "page": "page",
            "document": "document",
            "confidential": "confidentiel",
            "client": "client",
            "price": "prix",
            "product": "produit",
        },
    }
    dictionary = replacements.get(language.casefold(), {})
    translated_lines = []
    for line in text.splitlines():
        words = []
        for word in line.split():
            key = word.strip(".,;:!?()[]{}\"'").casefold()
            words.append(dictionary.get(key, word))
        translated_lines.append(" ".join(words))
    translated = "\n".join(translated_lines).strip()
    return (
        f"Traduction locale vers {language}\n"
        "Configurez OPENAI_API_KEY pour une traduction IA haute fidelite.\n\n"
        f"{translated}"
    )


def merge_pdfs(files: List) -> bytes:
    if len(files) < 2:
        raise ValueError("Au moins deux fichiers PDF sont necessaires pour fusionner.")

    writer = PdfWriter()
    temp_paths = []
    try:
        for uploaded in files:
            path = save_upload_file(uploaded)
            temp_paths.append(path)
            reader = _load_pdf_reader(path)
            for page in reader.pages:
                writer.add_page(page)
        if not writer.pages:
            raise ValueError("Aucune page PDF valide n'a ete trouvee.")
        return _serialize_writer(writer)
    finally:
        cleanup(temp_paths)


def split_pdf(file, pages: str) -> bytes:
    source_path = save_upload_file(file)
    try:
        reader = _load_pdf_reader(source_path)
        selected_pages = parse_page_ranges(pages, len(reader.pages))
        writer = PdfWriter()
        for page_index in selected_pages:
            writer.add_page(reader.pages[page_index])
        return _serialize_writer(writer)
    finally:
        cleanup([source_path])


def reorder_pdf_pages(file, pages: str) -> bytes:
    source_path = save_upload_file(file)
    try:
        reader = _load_pdf_reader(source_path)
        ordered_pages = parse_page_sequence(pages, len(reader.pages))
        writer = PdfWriter()
        for page_index in ordered_pages:
            writer.add_page(reader.pages[page_index])
        return _serialize_writer(writer)
    finally:
        cleanup([source_path])


def rotate_pdf(file, angle: int, pages: Optional[str]) -> bytes:
    source_path = save_upload_file(file)
    try:
        reader = _load_pdf_reader(source_path)
        writer = PdfWriter()
        total_pages = len(reader.pages)
        selected_pages = parse_page_ranges(pages, total_pages) if pages else list(range(total_pages))
        for page_index, page in enumerate(reader.pages):
            if page_index in selected_pages:
                page.rotate(angle)
            writer.add_page(page)
        return _serialize_writer(writer)
    finally:
        cleanup([source_path])


def crop_pdf(file, top: float = 0, right: float = 0, bottom: float = 0, left: float = 0) -> bytes:
    if min(top, right, bottom, left) < 0:
        raise ValueError("Les marges de recadrage doivent etre positives.")

    source_path = save_upload_file(file)
    try:
        reader = _load_pdf_reader(source_path)
        writer = PdfWriter()
        for page in reader.pages:
            width, height = _page_dimensions(page)
            if left + right >= width or top + bottom >= height:
                raise ValueError("Le recadrage est trop important pour au moins une page.")
            page.cropbox.lower_left = (float(page.mediabox.left) + left, float(page.mediabox.bottom) + bottom)
            page.cropbox.upper_right = (float(page.mediabox.right) - right, float(page.mediabox.top) - top)
            writer.add_page(page)
        return _serialize_writer(writer)
    finally:
        cleanup([source_path])


def compress_pdf(file) -> bytes:
    source_path = save_upload_file(file)
    compressed_path = Path(tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name)
    try:
        with pikepdf.open(source_path) as pdf:
            pdf.save(
                compressed_path,
                compress_streams=True,
                object_stream_mode=pikepdf.ObjectStreamMode.generate,
                recompress_flate=True,
            )
        return compressed_path.read_bytes()
    except pikepdf.PdfError as exc:
        raise ValueError("Impossible de compresser ce PDF.") from exc
    finally:
        cleanup([source_path, compressed_path])


def repair_pdf(file) -> bytes:
    source_path = save_upload_file(file)
    repaired_path = Path(tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name)
    try:
        with pikepdf.open(source_path, allow_overwriting_input=True) as pdf:
            pdf.save(repaired_path, fix_metadata_version=True, compress_streams=True)
        return repaired_path.read_bytes()
    except pikepdf.PdfError as exc:
        raise ValueError("Impossible de reparer ce PDF automatiquement.") from exc
    finally:
        cleanup([source_path, repaired_path])


def image_to_pdf(file) -> bytes:
    source_path = save_upload_file(file)
    try:
        image = Image.open(source_path)
        output = io.BytesIO()
        if image.mode in ("RGBA", "LA"):
            background = Image.new("RGB", image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[-1])
            image = background
        elif image.mode != "RGB":
            image = image.convert("RGB")
        image.save(output, format="PDF")
        output.seek(0)
        return output.read()
    except UnidentifiedImageError as exc:
        raise ValueError("Le fichier image est invalide ou n'est pas supporte.") from exc
    finally:
        cleanup([source_path])


def delete_pdf_pages(file, pages: str) -> bytes:
    source_path = save_upload_file(file)
    try:
        reader = _load_pdf_reader(source_path)
        total_pages = len(reader.pages)
        pages_to_delete = set(parse_page_ranges(pages, total_pages))
        if len(pages_to_delete) >= total_pages:
            raise ValueError("Impossible de supprimer toutes les pages du document.")
        writer = PdfWriter()
        for page_index, page in enumerate(reader.pages):
            if page_index not in pages_to_delete:
                writer.add_page(page)
        return _serialize_writer(writer)
    finally:
        cleanup([source_path])


def extract_pdf_pages(file, pages: str) -> bytes:
    return split_pdf(file, pages)


def add_watermark(file, text: str, opacity: float = 0.3) -> bytes:
    watermark_text = (text or "").strip()
    if not watermark_text:
        raise ValueError("Le texte du filigrane est obligatoire.")
    if not 0 < opacity <= 1:
        raise ValueError("L'opacite doit etre comprise entre 0 et 1.")

    source_path = save_upload_file(file)
    try:
        reader = _load_pdf_reader(source_path)
        writer = PdfWriter()
        for page in reader.pages:
            width, height = _page_dimensions(page)
            max_text_width = ((width ** 2 + height ** 2) ** 0.5) * 0.72
            font_size = _fit_font_size(
                watermark_text,
                "Helvetica-Bold",
                max_text_width,
                start_size=max(min(width, height) * 0.12, 28),
                min_size=18,
            )
            overlay = _build_overlay_page(
                width,
                height,
                lambda pdf_canvas, label=watermark_text, size=font_size: _draw_watermark(
                    pdf_canvas, width, height, label, size, opacity
                ),
            )
            page.merge_page(overlay)
            writer.add_page(page)
        return _serialize_writer(writer)
    finally:
        cleanup([source_path])


def pdf_to_images(file) -> bytes:
    source_path = save_upload_file(file)
    images: List[Image.Image] = []
    try:
        images = _render_pdf_pages(source_path)
        if not images:
            raise ValueError("Aucune page n'a pu etre convertie en image.")
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for page_number, image in enumerate(images, start=1):
                image_buffer = io.BytesIO()
                image.save(image_buffer, format="JPEG", quality=90)
                archive.writestr(f"page-{page_number}.jpg", image_buffer.getvalue())
        zip_buffer.seek(0)
        return zip_buffer.read()
    finally:
        for image in images:
            try:
                image.close()
            except Exception:
                pass
        cleanup([source_path])


def pdf_to_word(file) -> bytes:
    source_path = save_upload_file(file)
    try:
        page_texts = _extract_pdf_text_by_page(source_path)
        document = Document()
        document.add_heading("Export PDF vers Word", level=1)
        if not any(page_texts):
            raise ValueError("Aucun texte exploitable trouve dans ce PDF. OCR image indisponible ici.")
        for page_number, text in enumerate(page_texts, start=1):
            document.add_heading(f"Page {page_number}", level=2)
            document.add_paragraph(text or "[Page sans texte]")
        output = io.BytesIO()
        document.save(output)
        output.seek(0)
        return output.read()
    finally:
        cleanup([source_path])


def pdf_to_powerpoint(file) -> bytes:
    source_path = save_upload_file(file)
    images: List[Image.Image] = []
    try:
        images = _render_pdf_pages(source_path)
        if not images:
            raise ValueError("Aucune page n'a pu etre convertie en diapositives.")

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        blank_layout = presentation.slide_layouts[6]
        for image in images:
            slide = presentation.slides.add_slide(blank_layout)
            image_buffer = io.BytesIO()
            image.save(image_buffer, format="PNG")
            image_buffer.seek(0)
            slide.shapes.add_picture(
                image_buffer,
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )

        if len(presentation.slides) > len(images):
            presentation.slides._sldIdLst.remove(presentation.slides._sldIdLst[0])

        output = io.BytesIO()
        presentation.save(output)
        output.seek(0)
        return output.read()
    finally:
        for image in images:
            try:
                image.close()
            except Exception:
                pass
        cleanup([source_path])


def pdf_to_excel(file) -> bytes:
    source_path = save_upload_file(file)
    output = io.BytesIO()
    try:
        workbook = Workbook()
        default_sheet = workbook.active
        workbook.remove(default_sheet)

        with pdfplumber.open(str(source_path)) as pdf:
            created = False
            for page_number, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables() or []
                if tables:
                    for table_index, table in enumerate(tables, start=1):
                        sheet = workbook.create_sheet(title=f"P{page_number}T{table_index}")
                        for row_index, row in enumerate(table, start=1):
                            for col_index, cell in enumerate(row or [], start=1):
                                sheet.cell(row=row_index, column=col_index, value=(cell or "").strip())
                        created = True
                else:
                    text = (page.extract_text() or "").splitlines()
                    if text:
                        sheet = workbook.create_sheet(title=f"Page{page_number}")
                        for row_index, line in enumerate(text, start=1):
                            sheet.cell(row=row_index, column=1, value=line.strip())
                        created = True

        if not workbook.worksheets:
            raise ValueError("Aucune table ou texte exploitable trouve dans ce PDF.")

        workbook.save(output)
        output.seek(0)
        return output.read()
    finally:
        cleanup([source_path])


def extract_text_from_pdf(file) -> bytes:
    source_path = save_upload_file(file)
    try:
        ocr_text = _ocr_with_ocrmypdf(source_path, "fra+eng")
        if ocr_text:
            return ocr_text
        page_texts = _extract_pdf_text_by_page(source_path)
        if not any(page_texts):
            raise ValueError("Aucun texte exploitable trouve. OCR image indisponible dans cet environnement.")
        lines = []
        for page_number, text in enumerate(page_texts, start=1):
            lines.append(f"===== Page {page_number} =====")
            lines.append(text or "[Page sans texte]")
            lines.append("")
        return "\n".join(lines).encode("utf-8")
    finally:
        cleanup([source_path])


def censor_pdf(file, terms: str, case_sensitive: bool = False) -> bytes:
    source_path = save_upload_file(file)
    images: List[Image.Image] = []
    try:
        raw_terms = [term.strip() for term in terms.replace("\n", ",").split(",")]
        lookup_terms = [term for term in raw_terms if term]
        if not lookup_terms:
            raise ValueError("Indiquez au moins un terme a censurer.")

        images = _render_pdf_pages(source_path)
        if not images:
            raise ValueError("Impossible de rasteriser ce PDF pour la censure.")

        redacted_pages = 0
        with pdfplumber.open(str(source_path)) as pdf:
            for page_index, page in enumerate(pdf.pages):
                if page_index >= len(images):
                    break

                words = page.extract_words() or []
                image = images[page_index]
                draw = ImageDraw.Draw(image)
                scale_x = image.width / float(page.width or image.width)
                scale_y = image.height / float(page.height or image.height)
                page_redacted = False

                for word in words:
                    text = (word.get("text") or "").strip()
                    if not text:
                        continue

                    haystack = text if case_sensitive else text.casefold()
                    matches = any(
                        (term if case_sensitive else term.casefold()) in haystack
                        for term in lookup_terms
                    )
                    if not matches:
                        continue

                    x0 = max(0, int(float(word["x0"]) * scale_x) - 2)
                    x1 = min(image.width, int(float(word["x1"]) * scale_x) + 2)
                    top = max(0, int(float(word["top"]) * scale_y) - 2)
                    bottom = min(image.height, int(float(word["bottom"]) * scale_y) + 2)
                    draw.rectangle([x0, top, x1, bottom], fill="black")
                    page_redacted = True

                if page_redacted:
                    redacted_pages += 1

        if redacted_pages == 0:
            raise ValueError("Aucune occurrence des termes demandes n'a ete trouvee dans le PDF.")

        output = io.BytesIO()
        first, *rest = [image.convert("RGB") for image in images]
        first.save(output, format="PDF", save_all=True, append_images=rest, resolution=150.0)
        output.seek(0)
        return output.read()
    finally:
        for image in images:
            try:
                image.close()
            except Exception:
                pass
        cleanup([source_path])


def add_page_numbers(file, format_str: str = "{page}", position: str = "bottom-right") -> bytes:
    format_template = (format_str or "").strip()
    if not format_template:
        raise ValueError("Le format de numerotation est obligatoire.")
    if position not in PAGE_NUMBER_POSITIONS:
        raise ValueError("La position de numerotation demandee n'est pas supportee.")

    source_path = save_upload_file(file)
    try:
        reader = _load_pdf_reader(source_path)
        writer = PdfWriter()
        total_pages = len(reader.pages)
        for page_number, page in enumerate(reader.pages, start=1):
            try:
                label = format_template.format(page=page_number, total=total_pages)
            except (KeyError, ValueError) as exc:
                raise ValueError(
                    "Le format est invalide. Utilisez par exemple '{page}' ou 'Page {page}/{total}'."
                ) from exc

            width, height = _page_dimensions(page)
            font_name = "Helvetica"
            font_size = max(min(min(width, height) * 0.02, 14), 10)
            x, y, horizontal = _page_number_coordinates(width, height, position)
            overlay = _build_overlay_page(
                width,
                height,
                lambda pdf_canvas, text=label, draw_x=x, draw_y=y, size=font_size, align=horizontal: _draw_page_number(
                    pdf_canvas, text, draw_x, draw_y, font_name, size, align
                ),
            )
            page.merge_page(overlay)
            writer.add_page(page)
        return _serialize_writer(writer)
    finally:
        cleanup([source_path])


def sign_pdf(file, signer_name: str, reason: str = "", location: str = "", position: str = "bottom-right") -> bytes:
    signer = (signer_name or "").strip()
    if not signer:
        raise ValueError("Le nom du signataire est obligatoire.")
    if position not in SIGNATURE_POSITIONS:
        raise ValueError("La position de signature demandee n'est pas supportee.")

    source_path = save_upload_file(file)
    try:
        reader = _load_pdf_reader(source_path)
        writer = PdfWriter()
        signed_at = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
        lines = [
            f"Signe par {signer}",
            f"Date: {signed_at}",
        ]
        if reason.strip():
            lines.append(f"Motif: {reason.strip()}")
        if location.strip():
            lines.append(f"Lieu: {location.strip()}")
        lines.append("Signature visible non cryptographique")

        for page in reader.pages:
            width, height = _page_dimensions(page)
            overlay = _build_overlay_page(
                width,
                height,
                lambda pdf_canvas, draw_lines=lines: _draw_signature_box(
                    pdf_canvas, width, height, draw_lines, position
                ),
            )
            page.merge_page(overlay)
            writer.add_page(page)
        return _serialize_writer(writer)
    finally:
        cleanup([source_path])


def protect_pdf(file, user_password: str, owner_password: Optional[str] = None) -> bytes:
    if not user_password:
        raise ValueError("Le mot de passe utilisateur est obligatoire.")

    source_path = save_upload_file(file)
    protected_path = Path(tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name)
    try:
        with pikepdf.open(source_path) as pdf:
            pdf.save(
                protected_path,
                encryption=pikepdf.Encryption(
                    user=user_password,
                    owner=owner_password or user_password,
                    allow=pikepdf.Permissions(extract=False, modify_annotation=False),
                ),
            )
        return protected_path.read_bytes()
    except pikepdf.PdfError as exc:
        raise ValueError("Impossible de proteger ce PDF.") from exc
    finally:
        cleanup([source_path, protected_path])


def unlock_pdf(file, password: str) -> bytes:
    if not password:
        raise ValueError("Le mot de passe de deverrouillage est obligatoire.")

    source_path = save_upload_file(file)
    unlocked_path = Path(tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name)
    try:
        with pikepdf.open(source_path, password=password) as pdf:
            pdf.save(unlocked_path)
        return unlocked_path.read_bytes()
    except pikepdf.PasswordError as exc:
        raise ValueError("Mot de passe incorrect pour ce PDF.") from exc
    except pikepdf.PdfError as exc:
        raise ValueError("Impossible de deverrouiller ce PDF.") from exc
    finally:
        cleanup([source_path, unlocked_path])


def compare_pdfs(file_a, file_b) -> bytes:
    path_a = save_upload_file(file_a)
    path_b = save_upload_file(file_b)
    try:
        reader_a = _load_pdf_reader(path_a)
        reader_b = _load_pdf_reader(path_b)
        text_a = _extract_pdf_text_by_page(path_a)
        text_b = _extract_pdf_text_by_page(path_b)

        report = {
            "pages_a": len(reader_a.pages),
            "pages_b": len(reader_b.pages),
            "same_page_count": len(reader_a.pages) == len(reader_b.pages),
            "same_text": text_a == text_b,
            "page_differences": [],
        }

        for index in range(max(len(text_a), len(text_b))):
            page_a = text_a[index] if index < len(text_a) else ""
            page_b = text_b[index] if index < len(text_b) else ""
            if page_a != page_b:
                report["page_differences"].append(
                    {
                        "page": index + 1,
                        "text_a_preview": page_a[:200],
                        "text_b_preview": page_b[:200],
                    }
                )

        return json.dumps(report, ensure_ascii=False, indent=2).encode("utf-8")
    finally:
        cleanup([path_a, path_b])


def summarize_pdf(file, max_sentences: int = 6) -> bytes:
    if max_sentences < 1 or max_sentences > 20:
        raise ValueError("Le nombre de phrases du resume doit etre compris entre 1 et 20.")

    source_path = save_upload_file(file)
    try:
        text = _extract_full_pdf_text(source_path)
        ai_summary = _call_openai_text(
            "Tu resumes des documents PDF en francais, avec des points courts et fideles.",
            f"Resume ce document en {max_sentences} points maximum:\n\n{text}",
        )
        summary = ai_summary or _local_summary(text, max_sentences)
        return summary.encode("utf-8")
    finally:
        cleanup([source_path])


def translate_pdf(file, target_language: str = "francais") -> bytes:
    language = (target_language or "").strip()
    if not language:
        raise ValueError("La langue cible est obligatoire.")
    if len(language) > 40:
        raise ValueError("La langue cible est trop longue.")

    source_path = save_upload_file(file)
    try:
        text = _extract_full_pdf_text(source_path)
        ai_translation = _call_openai_text(
            "Tu traduis fidelement le texte fourni. Conserve les paragraphes et ne rajoute pas d'analyse.",
            f"Traduis ce texte vers {language}:\n\n{text}",
        )
        translated = ai_translation or _local_translate_text(text, language)
        return _text_lines_to_pdf(translated.splitlines(), title=f"Traduction vers {language}")
    finally:
        cleanup([source_path])


def html_to_pdf(file) -> bytes:
    source_path = save_upload_file(file)
    try:
        converted = _convert_with_libreoffice(source_path)
        if converted:
            return converted
        extractor = _HTMLTextExtractor()
        extractor.feed(source_path.read_text(encoding="utf-8", errors="ignore"))
        text = extractor.get_text()
        if not text:
            raise ValueError("Le fichier HTML ne contient pas de texte exploitable.")
        return _text_lines_to_pdf(text.splitlines(), title="Conversion HTML vers PDF")
    finally:
        cleanup([source_path])


def word_to_pdf(file) -> bytes:
    source_path = save_upload_file(file)
    try:
        converted = _convert_with_libreoffice(source_path)
        if converted:
            return converted
        document = Document(str(source_path))
        lines = []
        for paragraph in document.paragraphs:
            if paragraph.text.strip():
                lines.append(paragraph.text.strip())
        for table in document.tables:
            lines.append("")
            for row in table.rows:
                row_values = [cell.text.strip() for cell in row.cells]
                lines.append(" | ".join(value for value in row_values if value))
        if not lines:
            raise ValueError("Le fichier Word ne contient pas de contenu exploitable.")
        return _text_lines_to_pdf(lines, title="Conversion Word vers PDF")
    finally:
        cleanup([source_path])


def excel_to_pdf(file) -> bytes:
    source_path = save_upload_file(file)
    try:
        converted = _convert_with_libreoffice(source_path)
        if converted:
            return converted
        workbook = load_workbook(str(source_path), data_only=True)
        lines = []
        for worksheet in workbook.worksheets:
            lines.append(f"Feuille: {worksheet.title}")
            for row in worksheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell not in (None, "")]
                if values:
                    lines.append(" | ".join(values))
            lines.append("")
        if not _normalize_lines(lines):
            raise ValueError("Le fichier Excel ne contient pas de donnees exploitables.")
        return _text_lines_to_pdf(lines, title="Conversion Excel vers PDF")
    finally:
        cleanup([source_path])


def powerpoint_to_pdf(file) -> bytes:
    source_path = save_upload_file(file)
    try:
        converted = _convert_with_libreoffice(source_path)
        if converted:
            return converted
        presentation = Presentation(str(source_path))
        lines = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            lines.append(f"Diapositive {slide_number}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.extend(part.strip() for part in shape.text.splitlines() if part.strip())
            lines.append("")
        if not _normalize_lines(lines):
            raise ValueError("Le fichier PowerPoint ne contient pas de texte exploitable.")
        return _text_lines_to_pdf(lines, title="Conversion PowerPoint vers PDF")
    finally:
        cleanup([source_path])
